import os
import json
import requests
import time
import re
import random
import base64
import pandas as pd
from urllib.parse import quote
from dotenv import load_dotenv
from openai import OpenAI
from docx import Document
from docx.shared import Inches as DocxInches, Pt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF
import pypdf
import openpyxl
from openpyxl.chart import PieChart, BarChart, Reference

# Load environment variables from .env file
load_dotenv(override=True)

def clean_spacing(text):
    if not isinstance(text, str): 
        return text
    text = re.sub(r'(\d+)([a-zA-Z])', r'\1 \2', text)
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
    return re.sub(r' +', ' ', text).strip()

def encode_image_to_base64(file_path):
    with open(file_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def extract_file_text(file_path):
    if not file_path or not os.path.exists(file_path):
        return ""
    ext = os.path.splitext(file_path)[1].lower()
    text_content = ""
    try:
        if ext in ['.txt', '.md', '.json', '.csv']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
        elif ext == '.docx':
            doc = Document(file_path)
            text_content = "\n".join([p.text for p in doc.paragraphs if p.text])
        elif ext == '.pdf':
            reader = pypdf.PdfReader(file_path)
            text_content = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif ext in ['.xlsx', '.xls']:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            text_content = "\n".join([f"Sheet {name}:\n{df.to_csv(index=False)}" for name, df in df_dict.items()])
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            slides_text = []
            for s_idx, slide in enumerate(prs.slides):
                slide_words = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                if slide_words:
                    slides_text.append(f"Slide {s_idx+1}: " + " ".join(slide_words))
            text_content = "\n".join(slides_text)
        elif ext in ['.mp4', '.mov', '.avi', '.mkv']:
            text_content = f"[Reference Context: User uploaded a video file named '{os.path.basename(file_path)}']"
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return text_content[:6000]

class UniversalAgent:
    def __init__(self, base_url=None, api_key=None):
        server_url = base_url or os.getenv("LLAMA_SERVER_URL", "http://localhost:8080/v1")
        secret_key = api_key or os.getenv("LLAMA_API_KEY", "llama-cpp")
        
        self.client = OpenAI(base_url=server_url, api_key=secret_key)
        self.model = os.getenv("MODEL_NAME", "gemma-4-E4B-it")
        self.session = requests.Session()
        self.session.headers.update({
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        })

    def get_text(self, prompt, is_json=False):
        try:
            kwargs = {
                "model": self.model,
                "messages": [{"role": "user", "content": prompt}],
            }
            if is_json:
                kwargs["response_format"] = {"type": "json_object"}

            response = self.client.chat.completions.create(**kwargs)
            return response.choices[0].message.content
        except Exception as e:
            return "{}" if is_json else f"Error: {e}"

    def get_vision_response(self, prompt, image_path):
        try:
            if not image_path or not os.path.exists(image_path):
                return "❌ Error: Image file was not found on disk."

            base64_image = encode_image_to_base64(image_path)
            ext = os.path.splitext(image_path)[1].lower().replace('.', '')
            mime_type = "image/jpeg" if ext in ['jpg', 'jpeg'] else f"image/{ext}"
            
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url", 
                            "image_url": {
                                "url": f"data:{mime_type};base64,{base64_image}"
                            }
                        }
                    ]
                }]
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"❌ Vision Analysis Error: {e}"

    def generate_chat_title(self, first_prompt):
        try:
            p = f"Summarize this prompt into a clean 3-5 word title (NO quotes, NO punctuation): '{first_prompt}'"
            title = self.get_text(p).strip().strip('"').strip("'")
            return title if title and len(title) < 40 else first_prompt[:25] + "..."
        except Exception:
            return first_prompt[:25] + "..."

    def _extract_json(self, raw_text):
        cleaned = re.sub(r'```json\s*', '', raw_text)
        cleaned = re.sub(r'```\s*', '', cleaned).strip()
        match = re.search(r'[\{\[].*[\}\]]', cleaned, re.DOTALL)
        if match:
            cleaned = match.group(0)
        try:
            return json.loads(cleaned)
        except Exception:
            return {}

    def _fetch_stock_image(self, keyword, custom_filename=None):
        try:
            stop_words = {'overview', 'summary', 'introduction', 'part', 'slide', 'details', 'section', 'about', 'guide'}
            raw_words = re.sub(r'[^a-zA-Z0-9\s]', '', str(keyword)).split()
            clean_words = [w.lower() for w in raw_words if len(w) > 2 and w.lower() not in stop_words]
            search_term = clean_words[0] if clean_words else "technology"
            
            unique_seed = f"{int(time.time() * 1000)}_{random.randint(1000, 9999)}"
            url = f"https://loremflickr.com/800/600/{quote(search_term)}?random={unique_seed}"
            
            r = self.session.get(url, timeout=15, allow_redirects=True)
            if r.status_code == 200 and len(r.content) > 1000:
                filename = custom_filename if custom_filename else f"stock_{unique_seed}.jpg"
                with open(filename, "wb") as f:
                    f.write(r.content)
                return filename
            return None
        except Exception: 
            return None

    def _generate_ai_image(self, prompt_text, custom_filename=None, model="flux"):
        try:
            clean_prompt = quote(str(prompt_text).replace("\n", " ").strip() or "modern visual art")
            seed = random.randint(1, 2147483647)
            
            url = f"https://image.pollinations.ai/prompt/{clean_prompt}?seed={seed}&model={model}&width=1024&height=1024&nologo=true"
            
            res = self.session.get(url, timeout=45, allow_redirects=True)
            content_type = res.headers.get("Content-Type", "")
            
            if res.status_code == 200 and ("image" in content_type or len(res.content) > 5000):
                filename = custom_filename if custom_filename else f"ai_img_{seed}.png"
                if not filename.endswith(('.jpg', '.jpeg', '.png')):
                    filename += ".png"
                with open(filename, "wb") as f:
                    f.write(res.content)
                return filename
            else:
                return None
        except Exception:
            return None

    def _generate_ai_video(self, prompt_text, custom_filename=None, model="Lightricks/LTX-Video-0.9.8-13B-distilled"):
        try:
            hf_token = os.getenv("HF_TOKEN")
            if not hf_token:
                return None

            api_url = f"https://api-inference.huggingface.co/models/{model}"
            headers = {
                "Authorization": f"Bearer {hf_token}",
                "Content-Type": "application/json"
            }
            payload = {
                "inputs": str(prompt_text).strip(),
                "parameters": {"num_frames": 16}
            }

            res = requests.post(api_url, headers=headers, json=payload, timeout=120)
            if res.status_code == 200:
                seed = random.randint(1000, 9999)
                filename = custom_filename if custom_filename else f"ai_vid_{seed}.mp4"
                if not filename.endswith('.mp4'):
                    filename += '.mp4'
                with open(filename, "wb") as f:
                    f.write(res.content)
                return filename
            else:
                return None
        except Exception:
            return None

    def create_image(self, topic, filename):
        try:
            if not filename or filename == "output":
                filename = f"{topic.replace(' ', '_')[:25]}_{int(time.time())}.png"
            elif not filename.endswith(('.jpg', '.png', '.jpeg')):
                filename = f"{filename}.png"
                
            img_path = self._generate_ai_image(topic, custom_filename=filename)
            if img_path:
                return {"message": f"🖼️ AI Image Generated: {topic.title()}", "file_path": img_path}
            return {"message": "❌ Failed to generate AI image.", "file_path": None}
        except Exception as e:
            return {"message": f"❌ Image Error: {e}", "file_path": None}

    def create_video(self, topic, filename, model="Lightricks/LTX-Video-0.9.8-13B-distilled"):
        try:
            if not filename.endswith('.mp4'):
                filename = f"{topic.replace(' ', '_')[:25]}_{int(time.time())}.mp4"
            vid_path = self._generate_ai_video(topic, custom_filename=filename, model=model)
            if vid_path:
                return {"message": f"🎥 AI Video Generated: {topic.title()}", "file_path": vid_path}
            return {"message": "❌ Failed to generate AI video. Check HF_TOKEN environment variable.", "file_path": None}
        except Exception as e:
            return {"message": f"❌ Video Error: {e}", "file_path": None}

    def create_pdf(self, topic, filename, pages=1):
        try:
            formatted_topic = topic.title()
            if not filename.endswith(".pdf"): 
                filename = f"{topic.replace(' ', '_')[:25]}_{int(time.time())}.pdf"
            
            prompt = f"""
            Write an in-depth document about '{formatted_topic}' designed to span exactly {pages} pages.
            Return JSON ONLY matching:
            {{
                "pages": [
                    {{
                        "heading": "Descriptive Heading",
                        "paragraphs": ["Paragraph 1...", "Paragraph 2...", "Paragraph 3...", "Paragraph 4..."],
                        "stock_keyword": "visual noun"
                    }}
                ]
            }}
            Generate exactly {pages} sections.
            """
            
            raw_data = self.get_text(prompt, is_json=True)
            data = self._extract_json(raw_data)
            
            page_list = data.get("pages") or data.get("sections") if isinstance(data, dict) else []
            if not page_list:
                page_list = [{
                    "heading": f"{formatted_topic}: Architecture & Design",
                    "paragraphs": [
                        f"The development of modern {formatted_topic} represents a convergence of engineering and design principles.",
                        f"Analyzing system components highlights how balance across modules improves total operational reliability.",
                        f"From an implementation standpoint, deploying {formatted_topic} requires clear quality control standards.",
                        f"Future trajectories point toward integrated automation frameworks and optimized resource utilization."
                    ],
                    "stock_keyword": topic
                } for _ in range(pages)]

            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)

            for item in page_list:
                pdf.add_page()
                raw_heading = item.get("heading", formatted_topic)
                clean_heading = re.sub(r'(?i)\bpart\s*\d+\b|\bpage\s*\d+\b', '', raw_heading).strip(" :-") or formatted_topic
                
                pdf.set_font("helvetica", 'B', 15)
                pdf.cell(0, 8, txt=clean_spacing(clean_heading).title(), ln=True, align='L')
                pdf.ln(4)
                
                img = self._fetch_stock_image(item.get("stock_keyword", topic))
                if img and os.path.exists(img):
                    pdf.image(img, x=37.5, w=135)
                    pdf.ln(6)
                    os.remove(img)

                pdf.set_font("helvetica", size=10)
                for para in item.get("paragraphs", []):
                    clean_para = clean_spacing(para).encode('latin-1', 'ignore').decode('latin-1')
                    if clean_para:
                        pdf.multi_cell(0, 5.5, txt=clean_para)
                        pdf.ln(3.5)

            pdf.output(filename)
            return {"message": f"✅ PDF Created ({pdf.page} Pages): {formatted_topic}", "file_path": filename}
        except Exception as e: 
            return {"message": f"❌ PDF Error: {e}", "file_path": None}

    def create_word(self, topic, filename, pages=3):
        try:
            formatted_topic = topic.title()
            if not filename.endswith(".docx"): 
                filename = f"{topic.replace(' ', '_')[:25]}_{int(time.time())}.docx"
            
            prompt = f"""
            Write an in-depth technical report about '{formatted_topic}' spanning {pages} pages in Microsoft Word.
            Return JSON ONLY matching:
            {{
                "sections": [
                    {{
                        "heading": "Descriptive Heading",
                        "paragraphs": ["Paragraph 1...", "Paragraph 2...", "Paragraph 3...", "Paragraph 4..."],
                        "stock_keyword": "visual noun"
                    }}
                ]
            }}
            Generate exactly {pages} sections.
            """
            
            raw_data = self.get_text(prompt, is_json=True)
            data = self._extract_json(raw_data)
            
            section_list = data.get("sections") or data.get("pages") if isinstance(data, dict) else []
            if not section_list:
                section_list = [{
                    "heading": f"{formatted_topic}: Core Analysis",
                    "paragraphs": [
                        f"The history and engineering behind modern {formatted_topic} highlight essential industrial paradigms.",
                        f"Operational mechanisms rely on consistent design parameters to ensure high performance under strain.",
                        f"Industry adoption requires evaluating operational overhead, maintenance logistics, and scalability.",
                        f"Next-generation applications continue to incorporate automated analytics and adaptive frameworks."
                    ],
                    "stock_keyword": topic
                } for _ in range(pages)]

            doc = Document()
            for idx, sec in enumerate(section_list):
                if idx > 0:
                    doc.add_page_break()
                
                raw_heading = sec.get("heading", formatted_topic)
                clean_heading = re.sub(r'(?i)\bpart\s*\d+\b|\bpage\s*\d+\b', '', raw_heading).strip(" :-") or formatted_topic

                doc.add_heading(clean_spacing(clean_heading).title(), level=1)
                
                img_path = self._fetch_stock_image(sec.get("stock_keyword", topic))
                if img_path and os.path.exists(img_path):
                    doc.add_picture(img_path, width=DocxInches(5.0))
                    os.remove(img_path)
                
                for para in sec.get("paragraphs", []):
                    clean_para = clean_spacing(para)
                    if clean_para:
                        p = doc.add_paragraph(clean_para)
                        p.paragraph_format.space_after = Pt(8)
                        p.paragraph_format.line_spacing = 1.15
            
            doc.save(filename)
            return {"message": f"✅ Word Doc Created ({len(section_list)} Pages): {formatted_topic}", "file_path": filename}
        except Exception as e: 
            return {"message": f"❌ Word Error: {e}", "file_path": None}

    def create_excel(self, topic, filename, chart_type="none", reference_context="", count=30, original_file_path=None):
        try:
            formatted_topic = topic.title()
            
            if original_file_path and (filename.lower() == os.path.basename(original_file_path).lower() or os.path.abspath(filename) == os.path.abspath(original_file_path)):
                filename = f"updated_{os.path.basename(original_file_path)}"
            elif not filename.endswith(".xlsx"): 
                filename = f"{topic.replace(' ', '_')[:25]}_{int(time.time())}.xlsx"
            
            context_prompt = f"\nExisting Data from Uploaded File:\n{reference_context}\n" if reference_context else ""
            
            prompt = f"""
            {context_prompt}
            Generate a structured table for '{formatted_topic}'.
            
            CRITICAL REQUIREMENT:
            1. You MUST generate EXACTLY {count} distinct data rows in the "rows" array.
            2. Columns MUST be ["Name", "Subject", "Marks"].
            3. Keep the "Subject" column consistent (e.g. use '{formatted_topic}' or one single subject for all students).
            4. Do NOT stop after 2 or 3 rows. Include all {count} rows.

            Return JSON ONLY matching:
            {{
                "cols": ["Name", "Subject", "Marks"],
                "rows": [
                    ["Alex Johnson", "{formatted_topic}", 88],
                    ["Jordan Lee", "{formatted_topic}", 92]
                ],
                "should_chart": false
            }}
            """
            
            raw_data = self.get_text(prompt, is_json=True)
            data = self._extract_json(raw_data)
            
            cols = data.get("cols") if isinstance(data.get("cols"), list) and len(data.get("cols")) > 0 else ["Name", "Subject", "Marks"]
            rows = data.get("rows") if isinstance(data.get("rows"), list) else []

            if len(rows) < count:
                first_names = ["Alex", "Jordan", "Taylor", "Morgan", "Sam", "Chris", "Pat", "Riley", "Avery", "Dakota",
                               "Lucas", "Emma", "Liam", "Olivia", "Noah", "Ava", "Ethan", "Sophia", "Mason", "Isabella",
                               "Logan", "Mia", "James", "Harper", "Benjamin", "Evelyn", "Jacob", "Abigail", "Elijah", "Emily"]
                
                selected_subject = formatted_topic if formatted_topic and formatted_topic.lower() != "output" else "Mathematics"
                
                while len(rows) < count:
                    i = len(rows) + 1
                    name = f"{random.choice(first_names)} {chr(65 + (i % 26))}."
                    mark = random.randint(55, 98)
                    
                    if len(cols) >= 3:
                        rows.append([name, selected_subject, mark])
                    elif len(cols) == 2:
                        rows.append([name, mark])
                    else:
                        rows.append([name, selected_subject, mark])

            df = pd.DataFrame(rows, columns=cols[:len(rows[0])])
            
            for col in df.columns:
                try:
                    df[col] = pd.to_numeric(df[col])
                except (ValueError, TypeError):
                    pass
            
            with pd.ExcelWriter(filename, engine='openpyxl') as writer:
                df.to_excel(writer, sheet_name='Data', index=False)
            
            c_key = str(chart_type).lower().strip()
            should_chart = data.get("should_chart", False) or c_key in ["pie", "vertical_bar", "horizontal_bar", "bar", "column"]
            
            if should_chart and len(cols) >= 2:
                wb = openpyxl.load_workbook(filename)
                ws = wb['Data']
                
                if "vertical" in c_key or "col" in c_key:
                    chart = BarChart()
                    chart.type = "col"
                elif "horizontal" in c_key or "bar" in c_key:
                    chart = BarChart()
                    chart.type = "bar"
                else:
                    chart = PieChart()
                    
                chart.title = f"{formatted_topic} Overview"
                labels = Reference(ws, min_col=1, min_row=2, max_row=len(rows) + 1)
                chart_data = Reference(ws, min_col=2, min_row=1, max_row=len(rows) + 1)
                
                chart.add_data(chart_data, titles_from_data=True)
                chart.set_categories(labels)
                chart.width = 16
                chart.height = 10
                
                ws.add_chart(chart, "E2")
                wb.save(filename)
            
            return {"message": f"✅ Excel Spreadsheet Generated ({len(rows)} rows): {formatted_topic}", "file_path": filename}
        except Exception as e: 
            return {"message": f"❌ Excel Error: {e}", "file_path": None}

    def create_ppt(self, topic, filename, slides=3):
        try:
            formatted_topic = topic.title()
            if not filename.endswith(".pptx"): 
                filename = f"{topic.replace(' ', '_')[:25]}_{int(time.time())}.pptx"
            
            prompt = f"""
            Create a presentation about '{formatted_topic}'.
            Return JSON ONLY matching:
            {{
                "presentation_title": "{formatted_topic}",
                "slides": [
                    {{
                        "title": "Slide Title",
                        "bullets": ["Point 1", "Point 2", "Point 3"],
                        "stock_keyword": "visual noun"
                    }}
                ]
            }}
            Generate exactly {slides} slides.
            """
            
            raw_data = self.get_text(prompt, is_json=True)
            data = self._extract_json(raw_data)
            
            prs = Presentation()
            prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
            blank_layout = prs.slide_layouts[6]
            
            t_slide = prs.slides.add_slide(blank_layout)
            bg = t_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(18, 30, 49)
            bg.line.fill.background()
            
            tb = t_slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.5), Inches(2.5))
            p = tb.text_frame.paragraphs[0]
            p.text = clean_spacing(data.get("presentation_title", formatted_topic))
            p.font.size, p.font.bold = Pt(44), True
            p.font.color.rgb = RGBColor(255, 255, 255)

            for item in data.get("slides", []):
                slide = prs.slides.add_slide(blank_layout)
                hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
                hdr.fill.solid()
                hdr.fill.fore_color.rgb = RGBColor(18, 30, 49)
                hdr.line.fill.background()
                
                stb = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8))
                pst = stb.text_frame.paragraphs[0]
                pst.text = clean_spacing(item.get("title", "Overview")).title()
                pst.font.size, pst.font.bold = Pt(26), True
                pst.font.color.rgb = RGBColor(255, 255, 255)
                
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(6.8), Inches(5.3))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(245, 247, 250)
                card.line.fill.background()
                
                bb = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(6.2), Inches(4.7))
                tf_b = bb.text_frame
                tf_b.word_wrap = True
                
                for idx, bullet in enumerate(item.get("bullets", [])[:3]):
                    pb = tf_b.paragraphs[0] if idx == 0 else tf_b.add_paragraph()
                    pb.text = f"•  {clean_spacing(bullet.lstrip('-*• '))}"
                    pb.font.size = Pt(17)
                    pb.font.color.rgb = RGBColor(40, 40, 40)
                    pb.space_after = Pt(18)
                
                img_path = self._fetch_stock_image(item.get("stock_keyword", topic))
                if img_path and os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(8.0), Inches(1.6), width=Inches(4.5))
                    os.remove(img_path)
            
            prs.save(filename)
            return {"message": f"✅ PPT Created ({len(data.get('slides', []))} Slides): {formatted_topic}", "file_path": filename}
        except Exception as e:
            return {"message": f"❌ PPT Error: {e}", "file_path": None}

    def handle_request(self, user_prompt, file_path=None):
        reference_text = extract_file_text(file_path) if file_path else ""
        is_image_file = file_path and os.path.exists(file_path) and os.path.splitext(file_path)[1].lower() in ['.jpg', '.jpeg', '.png', '.webp']
        
        brain_p = f"""
        User Request: "{user_prompt}"
        Uploaded Context Available: {"Yes" if (reference_text or is_image_file) else "No"}
        
        Classify intent into ONE tool:
        - 'word': Multi-paragraph reports, written documents, articles, summaries.
        - 'pdf': Formatted multi-page PDF documents.
        - 'excel': Spreadsheets, data tables, lists, grades/marks, financial data.
        - 'ppt': Slide presentations.
        - 'image': EXCLUSIVELY when the user asks to CREATE, GENERATE, DRAW, or MAKE an image or picture.
        - 'video': EXCLUSIVELY when the user asks to CREATE, GENERATE, ANIMATE, or MAKE a video, animation, or clip.
        - 'text': Short direct conversational responses, Q&A, or analyzing uploaded files without generating a new media/doc file.

        Chart Types (For Excel ONLY if user explicitly asks for a visual graph/chart):
        - 'pie', 'vertical_bar', or 'horizontal_bar'
        - Default to 'none' if no chart was requested.

        Return JSON ONLY:
        {{
            "tool": "word | pdf | excel | ppt | image | video | text",
            "subject": "Core Subject",
            "file": "filename",
            "count": 3,
            "chart_type": "none | pie | vertical_bar | horizontal_bar"
        }}
        """
        try:
            raw_res = self.get_text(brain_p, is_json=True)
            res = self._extract_json(raw_res)
            
            t = res.get('tool', 'text')
            s = clean_spacing(res.get('subject', user_prompt))
            f = res.get('file', 'output')
            c_type = res.get('chart_type', 'none')
            
            num_match = re.search(r'\b(\d+)\b', user_prompt)
            if num_match:
                count = int(num_match.group(1))
            else:
                try:
                    count = int(res.get('count', 30))
                except (ValueError, TypeError):
                    count = 30

            count = min(max(count, 1), 100)
            full_prompt = f"{user_prompt}\n\nReference Context:\n{reference_text}" if reference_text else user_prompt

            if t == 'excel': 
                result = self.create_excel(
                    s, 
                    f, 
                    chart_type=c_type, 
                    reference_context=reference_text, 
                    count=count, 
                    original_file_path=file_path
                )
            elif t == 'word': 
                result = self.create_word(s, f, pages=count)
            elif t == 'pdf': 
                result = self.create_pdf(s, f, pages=count)
            elif t == 'ppt': 
                result = self.create_ppt(s, f, slides=count)
            elif t == 'image': 
                result = self.create_image(s, f)
            elif t == 'video':
                result = self.create_video(s, f)
            else: 
                if is_image_file:
                    text_out = self.get_vision_response(user_prompt, file_path)
                else:
                    text_out = self.get_text(full_prompt)
                result = {"message": text_out, "file_path": None}

            return json.dumps(result)
            
        except Exception as e:
            return json.dumps({"message": f"❌ Error: {e}", "file_path": None})